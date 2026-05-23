"""Génère SOUTENANCE.pptx — présentation ~10 slides pour la soutenance ESSAI 1A.

Sortie : Projet/SOUTENANCE.pptx
Palette inspirée du dashboard-v2 (cosmic climate).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent.parent
OUTPATH = HERE / "SOUTENANCE.pptx"

# ─────────── Palette cosmic climate ───────────
BG_DARK   = RGBColor(0x0A, 0x0E, 0x1A)   # deep space
BG_PANEL  = RGBColor(0x14, 0x1C, 0x2D)   # slightly lighter panel
TXT       = RGBColor(0xF5, 0xF5, 0xF7)   # near-white
TXT_DIM   = RGBColor(0x9C, 0xA3, 0xAF)   # dim grey
HOT1      = RGBColor(0xFF, 0x6B, 0x35)   # warming orange
HOT2      = RGBColor(0xF7, 0x93, 0x1E)
HOT3      = RGBColor(0xFF, 0xB6, 0x27)
COLD1     = RGBColor(0x00, 0xD9, 0xFF)   # CO2 cyan
COLD2     = RGBColor(0x00, 0x77, 0xB6)
GREEN     = RGBColor(0x52, 0xFF, 0xB8)   # biosphere
ACCENT    = RGBColor(0x4D, 0xAB, 0xFF)

# ─────────── Layout helpers ───────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_blank(prs):
    """Add a blank slide with dark background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return slide


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=TXT, font="Inter",
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_kicker(slide, text, color=COLD1):
    """Section kicker top-left."""
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.4),
                text, size=11, bold=True, color=color, font="Consolas")


def add_section_label(slide, num, label):
    """Top: `01 · LABEL`"""
    txt = f"{num}  ·  {label.upper()}"
    add_kicker(slide, txt)


def add_title(slide, text, size=40, color=TXT, top=Inches(0.85)):
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(1.2),
                text, size=size, bold=True, color=color)


def add_bullets(slide, items, *, left, top, width, height,
                size=15, color=TXT, line_gap=8):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, c = item
        else:
            text, c = item, color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(line_gap)
        run = p.add_run()
        run.text = "▸ " + text
        run.font.name = "Inter"
        run.font.size = Pt(size)
        run.font.color.rgb = c
    return tb


def add_panel(slide, left, top, width, height, color=BG_PANEL):
    """Glass-morphism panel (filled rect)."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.06
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.color.rgb = RGBColor(0x2A, 0x36, 0x4E)
    rect.line.width = Pt(0.75)
    return rect


def add_kpi(slide, left, top, w, h, value, label, *, accent=HOT1):
    add_panel(slide, left, top, w, h)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1), w, Inches(0.35),
                value, size=28, bold=True, color=accent)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.6), w, Inches(0.3),
                label, size=10, color=TXT_DIM)


def add_image(slide, path, left, top, width=None, height=None):
    if not path.exists():
        return None
    if width:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def add_footer(slide, page, total=10):
    add_textbox(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
                "ESSAI 1A · Programmation Mathématique 2025–2026",
                size=9, color=TXT_DIM, font="Consolas")
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                f"{page:02d} / {total:02d}",
                size=9, color=TXT_DIM, font="Consolas", align=PP_ALIGN.RIGHT)


# ─────────── Build presentation ───────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ═══════════════════════ SLIDE 1 — TITRE ═══════════════════════
s = new_blank(prs)
# Glow accent line
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(0.6), Inches(2.4), Inches(0.5), Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = HOT1; line.line.fill.background()

add_textbox(s, Inches(0.6), Inches(1.5), Inches(8), Inches(0.5),
            "CLIMAT  ·  CO₂  ·  1979–2025", size=14, color=COLD1, font="Consolas")

add_textbox(s, Inches(0.6), Inches(2.7), Inches(12), Inches(1.5),
            "Migration R → Python", size=54, bold=True, color=TXT)

add_textbox(s, Inches(0.6), Inches(3.7), Inches(12), Inches(1.5),
            "+ Tableau de bord interactif", size=54, bold=True, color=HOT1)

add_textbox(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.6),
            "Quel est le lien statistique entre 18 variables climatiques NOAA",
            size=16, color=TXT_DIM)
add_textbox(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.6),
            "et la concentration atmosphérique de CO₂ ?",
            size=16, color=TXT_DIM)

add_textbox(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
            "ESSAI 1A — Programmation Mathématique — Mai 2026",
            size=11, color=TXT_DIM, font="Consolas")

# ═══════════════════════ SLIDE 2 — SOURCES DE DONNÉES ═══════════════════════
s = new_blank(prs)
add_section_label(s, "01", "Sources & formats")
add_title(s, "8 sources, 4 formats, 60 Go de données brutes.")

headers = ["Source", "Période", "Résolution", "Format", "Volume"]
rows = [
    ("CO₂ mensuel global · NOAA GML",    "1979–2025", "mensuel",       "CSV",   "24 Ko"),
    ("CO₂ Mauna Loa · NOAA/SIO",          "1958–2026", "mensuel",       "CSV",   "23 Ko"),
    ("CO₂ paléo · Vostok",                "414 ka BP",  "irrégulier n=283", "CSV", "8 Ko"),
    ("Émissions fossiles · GCB 2025",     "1750–2024", "annuel",        "CSV",   "35 Ko"),
    ("ENSO ONI · NOAA CPC",               "1950–pres.","trimestriel",   "CSV",   "19 Ko"),
    ("Réanalyse CFSR/CFSv2 · 2.5°",       "1979–2026", "144 × 73",      "GRIB2", "3.5 Go"),
    ("Réanalyse CFSR/CFSv2 · 0.5°",       "1979–2026", "720 × 361",     "GRIB2", "57 Go"),
]

# Build table
left = Inches(0.6); top = Inches(2.2)
col_widths = [Inches(4.4), Inches(1.8), Inches(2.0), Inches(1.4), Inches(1.4)]
total_w = sum(c for c in col_widths)
row_h = Inches(0.42)

# Header
x = left
for h, cw in zip(headers, col_widths):
    cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, cw, row_h)
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1C, 0x2A, 0x44)
    cell.line.color.rgb = RGBColor(0x2A, 0x36, 0x4E); cell.line.width = Pt(0.5)
    tb = cell.text_frame; tb.margin_left = Inches(0.1); tb.margin_top = Inches(0.06)
    p = tb.paragraphs[0]; r = p.add_run(); r.text = h
    r.font.name = "Consolas"; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = COLD1
    x += cw

# Body
for i, row in enumerate(rows):
    x = left; y = top + row_h * (i + 1)
    bg = BG_PANEL if i % 2 == 0 else RGBColor(0x10, 0x18, 0x28)
    for txt, cw in zip(row, col_widths):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, row_h)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        cell.line.color.rgb = RGBColor(0x2A, 0x36, 0x4E); cell.line.width = Pt(0.3)
        tb = cell.text_frame; tb.margin_left = Inches(0.1); tb.margin_top = Inches(0.05)
        p = tb.paragraphs[0]; r = p.add_run(); r.text = txt
        r.font.name = "Inter"; r.font.size = Pt(11)
        r.font.color.rgb = TXT
        x += cw

# Note formats
add_textbox(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
            "▸  GRIB2 lu via wgrib2 → NetCDF → xarray  ·  CSV chargés via pandas  ·  0 NA après filtrage sur les 18 variables × 564 mois",
            size=12, color=TXT_DIM)

add_footer(s, 2)

# ═══════════════════════ SLIDE 3 — VARIABLES ═══════════════════════
s = new_blank(prs)
add_section_label(s, "02", "Variables retenues")
add_title(s, "18 variables CFSR  +  3 indicateurs dérivés CRE.")

# Left column: 18 base vars
vars_block = [
    "T2m, T500              température (2 m, 500 hPa)",
    "SPFH2m                 humidité spécifique 2 m",
    "PWAT                   eau précipitable colonne",
    "APCP                   précipitations cumulées",
    "TCDC                   couverture nuageuse",
    "DLWRF / ULWRF          LW descendant / ascendant",
    "DSWRF / USWRF          SW descendant / réfléchi",
    "PRMSL                  pression réduite mer",
    "CSDSF / CSUSF          SW clear-sky desc. / asc.",
    "CSDLF / CSULF          LW clear-sky desc. / asc.",
    "CDUVB / DUVB           UV-B clear-sky / all-sky",
    "ALBDO                  albédo de surface",
]
add_panel(s, Inches(0.6), Inches(2.1), Inches(7.4), Inches(4.6))
add_textbox(s, Inches(0.85), Inches(2.25), Inches(7), Inches(0.4),
            "▸  18 VARIABLES BRUTES  ·  GRIB2 → NetCDF", size=11, bold=True,
            color=COLD1, font="Consolas")
add_textbox(s, Inches(0.85), Inches(2.65), Inches(7), Inches(4),
            "\n".join(vars_block), size=11, color=TXT, font="Consolas")

# Right: dérivés + formules
add_panel(s, Inches(8.3), Inches(2.1), Inches(4.5), Inches(4.6))
add_textbox(s, Inches(8.55), Inches(2.25), Inches(4), Inches(0.4),
            "▸  3 INDICATEURS DÉRIVÉS  (CRE)", size=11, bold=True,
            color=HOT1, font="Consolas")
add_textbox(s, Inches(8.55), Inches(2.7), Inches(4), Inches(3),
            "Cloud Radiative Effect :", size=14, color=TXT, bold=True)
add_textbox(s, Inches(8.55), Inches(3.1), Inches(4), Inches(0.4),
            "CRE_SW  =  (DSWRF−USWRF)\n                 −(CSDSF−CSUSF)",
            size=11, color=GREEN, font="Consolas")
add_textbox(s, Inches(8.55), Inches(3.9), Inches(4), Inches(0.4),
            "CRE_LW  =  (DLWRF−ULWRF)\n                 −(CSDLF−CSULF)",
            size=11, color=GREEN, font="Consolas")
add_textbox(s, Inches(8.55), Inches(4.7), Inches(4), Inches(0.4),
            "CRE_net =  CRE_SW + CRE_LW",
            size=11, color=GREEN, font="Consolas")
add_textbox(s, Inches(8.55), Inches(5.5), Inches(4), Inches(1.2),
            "→ permet d'isoler l'effet des nuages\n     hors signal clear-sky",
            size=11, color=TXT_DIM)

add_footer(s, 3)

# ═══════════════════════ SLIDE 4 — PIPELINE R → PYTHON ═══════════════════════
s = new_blank(prs)
add_section_label(s, "03", "Migration R → Python")
add_title(s, "Tout le pipeline traduit, validé à 10⁻¹³ près.")

# Two-column table: R → Python
add_panel(s, Inches(0.6), Inches(2.1), Inches(8.5), Inches(4.4))
add_textbox(s, Inches(0.85), Inches(2.25), Inches(8), Inches(0.4),
            "▸  ÉQUIVALENCES TECHNIQUES", size=11, bold=True,
            color=COLD1, font="Consolas")

pipeline_rows = [
    ("dplyr / tidyr",                "→",  "pandas  (groupby, pivot, melt)"),
    ("ggplot2",                       "→",  "matplotlib + seaborn"),
    ("tsibble / forecast",            "→",  "statsmodels (ARIMA, STL, MK)"),
    ("ncdf4",                         "→",  "xarray + netCDF4"),
    ("Boxplot / IQR / z-score",       "→",  "scipy.stats + pandas"),
    ("Mann-Kendall + Sen + bootstrap","→",  "pymannkendall + np.random"),
    ("Granger causality",             "→",  "statsmodels.tsa.grangercausalitytests"),
    ("Régression stepwise + Lasso",   "→",  "sklearn.linear_model.LassoCV"),
]
y = Inches(2.75)
for r_side, arrow, py_side in pipeline_rows:
    add_textbox(s, Inches(0.85), y, Inches(3.6), Inches(0.4),
                r_side, size=12, color=TXT, font="Consolas")
    add_textbox(s, Inches(4.55), y, Inches(0.4), Inches(0.4),
                arrow, size=12, color=HOT1, font="Consolas")
    add_textbox(s, Inches(4.95), y, Inches(4.0), Inches(0.4),
                py_side, size=12, color=GREEN, font="Consolas")
    y += Inches(0.42)

# Right: validation
add_panel(s, Inches(9.4), Inches(2.1), Inches(3.4), Inches(4.4))
add_textbox(s, Inches(9.6), Inches(2.25), Inches(3), Inches(0.4),
            "▸  VALIDATION", size=11, bold=True, color=COLD1, font="Consolas")
add_textbox(s, Inches(9.6), Inches(2.7), Inches(3.1), Inches(0.6),
            "10⁻¹³", size=44, bold=True, color=GREEN)
add_textbox(s, Inches(9.6), Inches(3.5), Inches(3.1), Inches(0.5),
            "précision relative max", size=11, color=TXT_DIM)

add_textbox(s, Inches(9.6), Inches(4.2), Inches(3.1), Inches(2.2),
            "▸ Moyennes pondérées cos(lat)\n▸ Pentes Sen / IC bootstrap\n▸ Granger / Lasso / ARIMA\n▸ Cartes pixel 0.5°\n\nDifférence inhérente aux\nRNG sur bootstrap (~10⁻¹⁰)",
            size=11, color=TXT)

add_footer(s, 4)

# ═══════════════════════ SLIDE 5 — PHASE 1 CO₂ ═══════════════════════
s = new_blank(prs)
add_section_label(s, "04", "Phase 1 — Caractérisation du CO₂")
add_title(s, "Une hausse quasi-linéaire, 30× plus rapide que paléo.")

# KPI strip
add_kpi(s, Inches(0.6), Inches(2.1), Inches(2.4), Inches(1.0),
        "+1.881", "Pente Sen (ppm/an)", accent=HOT2)
add_kpi(s, Inches(3.2), Inches(2.1), Inches(2.4), Inches(1.0),
        "+89", "ppm cumulés sur 47 ans", accent=HOT1)
add_kpi(s, Inches(5.8), Inches(2.1), Inches(2.4), Inches(1.0),
        "30×", "vs paléo (Vostok)", accent=GREEN)
add_kpi(s, Inches(8.4), Inches(2.1), Inches(2.4), Inches(1.0),
        "[1.78, 1.97]", "IC95 bootstrap n=1000", accent=COLD1)

# Methods bullets
add_bullets(s, [
    "STL · décomposition saisonnière + tendance robuste (Loess)",
    "Mann-Kendall non-paramétrique  ·  τ ≈ 0.93,  p < 10⁻⁵⁰",
    "ARIMA(1,1,1)(1,1,1)₁₂  ·  prévision out-of-sample sur 5 ans",
    "Détection de ruptures Bayes (PELT)  ·  aucune rupture significative depuis 1979",
    "Comparaison Vostok (paléo) vs moderne  ·  taux ~0.06 ppm/an au pic naturel",
    "Anomalie COVID-19 modérée (-0.4 ppm cumulé)  ·  fraction airborne ≈ 53 % stable",
], left=Inches(0.6), top=Inches(3.5), width=Inches(12), height=Inches(3.4), size=14)

add_footer(s, 5)

# ═══════════════════════ SLIDE 6 — PHASE 2 CLIMAT 2.5° ═══════════════════════
s = new_blank(prs)
add_section_label(s, "05", "Phase 2 — Climat global 2.5°")
add_title(s, "5 représentations  ·  15/21 variables Granger-causales.", size=34)

# KPI strip
add_kpi(s, Inches(0.6), Inches(1.95), Inches(2.5), Inches(1.0),
        "15 / 21", "Granger X→CO₂ (d12)", accent=GREEN)
add_kpi(s, Inches(3.3), Inches(1.95), Inches(2.5), Inches(1.0),
        "R² = 0.748", "modèle multivar. Newey-West", accent=COLD1)
add_kpi(s, Inches(6.0), Inches(1.95), Inches(2.5), Inches(1.0),
        "+7.84", "W/m² · CSDLF · 47 ans", accent=HOT1)
add_kpi(s, Inches(8.7), Inches(1.95), Inches(2.5), Inches(1.0),
        "11/21", "spurious trend détecté", accent=HOT3)

# Body: 2 cols
add_panel(s, Inches(0.6), Inches(3.2), Inches(6.0), Inches(3.6))
add_textbox(s, Inches(0.85), Inches(3.35), Inches(5.8), Inches(0.4),
            "▸  TOP 5 CORR. RÉSIDUS (1979–2025)", size=11, bold=True,
            color=COLD1, font="Consolas")
rows_top = [
    ("CRE_LW",   "−0.523", HOT1),
    ("CRE_SW",   "+0.489", COLD1),
    ("DSWRF",    "+0.437", COLD1),
    ("CRE_net",  "+0.432", COLD1),
    ("PRMSL",    "−0.394", HOT1),
]
y = Inches(3.85)
for name, val, col in rows_top:
    add_textbox(s, Inches(0.95), y, Inches(3), Inches(0.4),
                name, size=14, color=TXT, font="Consolas")
    add_textbox(s, Inches(3.95), y, Inches(2), Inches(0.4),
                val, size=14, color=col, font="Consolas", bold=True)
    y += Inches(0.5)

add_panel(s, Inches(6.8), Inches(3.2), Inches(6.0), Inches(3.6))
add_textbox(s, Inches(7.05), Inches(3.35), Inches(5.8), Inches(0.4),
            "▸  5 REPRÉSENTATIONS TEMPORELLES", size=11, bold=True,
            color=HOT1, font="Consolas")
add_textbox(s, Inches(7.05), Inches(3.85), Inches(5.6), Inches(2.7),
            "level    →   série brute (signal séculaire)\n"
            "anom    →   anomalies désaisonnées\n"
            "resid    →   anomalies dé-trendées\n"
            "d1        →   différence 1er ordre  (mois-à-mois)\n"
            "d12      →   différence saisonnière (interannuel)\n\n"
            "→ révèle que 11/21 corrélations brutes\n"
            "    sont portées par la trend commune.",
            size=12, color=TXT, font="Consolas")

add_footer(s, 6)

# ═══════════════════════ SLIDE 7 — PHASE 3 CLIMAT 0.5° ═══════════════════════
s = new_blank(prs)
add_section_label(s, "06", "Phase 3 — Climat régional 0.5°")
add_title(s, "Cartes pixel par pixel  ·  4 hotspots  ·  amplification arctique 4×.", size=30)

# KPI
add_kpi(s, Inches(0.6), Inches(2.1), Inches(2.5), Inches(1.0),
        "259 920", "pixels × 564 mois", accent=COLD1)
add_kpi(s, Inches(3.3), Inches(2.1), Inches(2.5), Inches(1.0),
        "4.0×", "amplif. arctique vs global", accent=HOT1)
add_kpi(s, Inches(6.0), Inches(2.1), Inches(2.5), Inches(1.0),
        "0.999", "corr. 0.5° vs 2.5° (validation)", accent=GREEN)

# Bullets
add_bullets(s, [
    "5 bandes de latitude  ·  moyennes pondérées cos(lat) par pixel",
    "Pente Sen + MK p-value en pixel-par-pixel  (Sen vectorisé, MK sur 1/16 pixels)",
    "4 hotspots : Sahel · Sibérie centrale · Amazonie · Indonésie",
    "Régression multivariée + Granger par zone  (10 zones)",
    "R² par zone : global 0.75 · tropical 0.69 · boreal 0.28 · Amazonie 0.04",
    "Asymétrie hémisphérique :  T2m N/S = 1.98×  ·  PWAT N/S = 4.27×",
], left=Inches(0.6), top=Inches(3.4), width=Inches(12), height=Inches(3.4), size=14)

add_footer(s, 7)

# ═══════════════════════ SLIDE 8 — LIMITE CFSR/CFSv2 ═══════════════════════
s = new_blank(prs)
add_section_label(s, "07", "Regard critique")
add_title(s, "Le saut CFSR → CFSv2 (jan 2011) explique un tiers du R².",
          color=HOT1, size=30)

# Before/after R²
add_panel(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(2.0))
add_textbox(s, Inches(0.85), Inches(2.25), Inches(5.5), Inches(0.4),
            "▸  AVANT HOMOGÉNÉISATION", size=11, bold=True,
            color=HOT1, font="Consolas")
add_textbox(s, Inches(0.85), Inches(2.7), Inches(5), Inches(0.7),
            "R² = 0.748", size=40, bold=True, color=TXT)
add_textbox(s, Inches(0.85), Inches(3.5), Inches(5), Inches(0.5),
            "12 variables résiduelles, série brute", size=12, color=TXT_DIM)

add_panel(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(2.0))
add_textbox(s, Inches(7.05), Inches(2.25), Inches(5.5), Inches(0.4),
            "▸  APRÈS HOMOGÉNÉISATION (jump corrigé)", size=11, bold=True,
            color=GREEN, font="Consolas")
add_textbox(s, Inches(7.05), Inches(2.7), Inches(5), Inches(0.7),
            "R² = 0.46", size=40, bold=True, color=GREEN)
add_textbox(s, Inches(7.05), Inches(3.5), Inches(5), Inches(0.5),
            "≈ 1/3 du R² était artefact technique", size=12, color=TXT_DIM)

# Top jumps
add_panel(s, Inches(0.6), Inches(4.35), Inches(12.2), Inches(2.4))
add_textbox(s, Inches(0.85), Inches(4.45), Inches(11), Inches(0.4),
            "▸  17 / 21 VARIABLES AFFECTÉES (p < 0.05)   ·   TOP 5 JUMP_IN_SD",
            size=11, bold=True, color=HOT1, font="Consolas")
jumps = [
    ("CRE_LW",   "−1.92 sd",  "−8.0 %"),
    ("PRMSL",    "−1.77 sd",  "−0.03 %"),
    ("CRE_SW",   "+1.56 sd",  "−11.4 %"),
    ("TCDC",     "+1.14 sd",  "+4.3 %"),
    ("CRE_net",  "+1.13 sd",  "−15.7 %"),
]
y = Inches(4.95)
x0 = Inches(0.95)
add_textbox(s, x0, y, Inches(2.5), Inches(0.3), "Variable", size=10, bold=True,
            color=COLD1, font="Consolas")
add_textbox(s, x0 + Inches(3), y, Inches(2.5), Inches(0.3), "Jump (sd)",
            size=10, bold=True, color=COLD1, font="Consolas")
add_textbox(s, x0 + Inches(6), y, Inches(2.5), Inches(0.3), "Jump (%)",
            size=10, bold=True, color=COLD1, font="Consolas")
y += Inches(0.35)
for v, sd, pct in jumps:
    add_textbox(s, x0, y, Inches(2.5), Inches(0.3), v, size=12, color=TXT, font="Consolas")
    add_textbox(s, x0 + Inches(3), y, Inches(2.5), Inches(0.3), sd,
                size=12, color=HOT1, font="Consolas")
    add_textbox(s, x0 + Inches(6), y, Inches(2.5), Inches(0.3), pct,
                size=12, color=TXT_DIM, font="Consolas")
    y += Inches(0.28)

add_footer(s, 8)

# ═══════════════════════ SLIDE 9 — DASHBOARD ═══════════════════════
s = new_blank(prs)
add_section_label(s, "08", "Dashboard interactif")
add_title(s, "6 sections cinématiques  ·  globe 3D NASA + 36 grilles réelles.")

# Two columns: tech + sections
add_panel(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.6))
add_textbox(s, Inches(0.85), Inches(2.25), Inches(5.5), Inches(0.4),
            "▸  STACK TECHNIQUE", size=11, bold=True, color=COLD1, font="Consolas")
add_textbox(s, Inches(0.85), Inches(2.7), Inches(5.5), Inches(3.5),
            "React 18.3  +  Three.js 0.158  +  Babel in-browser\n"
            "Pas de bundler  ·  édition .jsx live  ·  cache-bust ?v=…\n\n"
            "Globe shader-based  +  texture NASA Blue Marble\n"
            "TopoJSON coastlines  ·  starfield  ·  hotspots interactifs\n\n"
            "Données scientifiques inline  +  fetch JSON 190 KB\n"
            "(36 grilles 18 vars × 2 modes  ·  export depuis .pkl)\n\n"
            "Server :  python -m http.server 8088",
            size=12, color=TXT, font="Consolas")

add_panel(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(4.6))
add_textbox(s, Inches(7.05), Inches(2.25), Inches(5.5), Inches(0.4),
            "▸  6 SECTIONS CINÉMATIQUES", size=11, bold=True,
            color=HOT1, font="Consolas")
sections_list = [
    ("01", "Hero", "+89 ppm en 47 ans — KPI vivants"),
    ("02", "Trajectoire", "Courbe CO₂ + 3 modèles + timeline événements"),
    ("03", "Earth 3D", "Globe interactif · 18 vars × 2 modes"),
    ("04", "Lien", "Heatmap 21×5 · Sankey Granger · R² 0.75"),
    ("05", "Hotspots", "4 mini-globes · amplif. arctique · radar"),
    ("06", "Critique", "Saut CFSR · 3 limites · sources"),
]
y = Inches(2.7)
for num, name, desc in sections_list:
    add_textbox(s, Inches(7.05), y, Inches(0.5), Inches(0.3),
                num, size=12, color=HOT3, font="Consolas", bold=True)
    add_textbox(s, Inches(7.55), y, Inches(2), Inches(0.3),
                name, size=12, color=TXT, font="Consolas", bold=True)
    add_textbox(s, Inches(7.05), y + Inches(0.25), Inches(5.5), Inches(0.3),
                desc, size=10, color=TXT_DIM)
    y += Inches(0.62)

add_footer(s, 9)

# ═══════════════════════ SLIDE 10 — CONCLUSIONS ═══════════════════════
s = new_blank(prs)
add_section_label(s, "09", "Conclusions")
add_title(s, "Trois résultats robustes  ·  honnêtes sur les limites.")

# 3 result cards
def result_card(left, top, w, h, num, title, body, accent):
    add_panel(s, left, top, w, h)
    add_textbox(s, left + Inches(0.2), top + Inches(0.15), w, Inches(0.5),
                num, size=24, bold=True, color=accent, font="Consolas")
    add_textbox(s, left + Inches(0.2), top + Inches(0.7), w, Inches(0.5),
                title, size=14, bold=True, color=TXT)
    add_textbox(s, left + Inches(0.2), top + Inches(1.15), w - Inches(0.3),
                h - Inches(1.2), body, size=11, color=TXT_DIM)

result_card(Inches(0.6), Inches(2.1), Inches(4.0), Inches(2.4),
            "01", "Corrélations brutes trompeuses",
            "11 variables sur 21 ont |r| > 0.4 sur les niveaux\n"
            "mais perdent ce lien sur les résidus.\n"
            "→ partage de la trend séculaire commune.",
            COLD1)
result_card(Inches(4.8), Inches(2.1), Inches(4.0), Inches(2.4),
            "02", "5 variables réellement actives",
            "À l'échelle interannuelle (résid + d12) :\n"
            "CRE_LW · CRE_SW · DSWRF · CRE_net · PRMSL.\n"
            "→ 15/21 causent Granger le CO₂ (lag 6 mois).",
            GREEN)
result_card(Inches(9.0), Inches(2.1), Inches(3.8), Inches(2.4),
            "03", "Une partie est artefact",
            "Saut CFSR → CFSv2 en jan 2011 :\n"
            "17/21 vars affectées.\n"
            "R² 0.75 → 0.46 après homogénéisation.\n"
            "T2m / CSDLF gagnent (forçage GES démasqué).",
            HOT1)

# Limites + perspectives
add_panel(s, Inches(0.6), Inches(4.7), Inches(12.2), Inches(2.0))
add_textbox(s, Inches(0.85), Inches(4.85), Inches(11), Inches(0.4),
            "▸  LIMITES ASSUMÉES  &  PERSPECTIVES", size=11, bold=True,
            color=HOT3, font="Consolas")
add_textbox(s, Inches(0.85), Inches(5.3), Inches(11.5), Inches(1.4),
            "▸ Détendrage linéaire alors que la trend CO₂ est cubique  →  GAM / splines.\n"
            "▸ Moyennes globales 2.5° masquent les puits tropicaux  →  c'est ce que résoud la Phase 3.\n"
            "▸ MK p-value en 0.5° calculé sur 1/16 pixels (coût)  ·  Sen complet vectorisé.\n"
            "▸ Migration validée 10⁻¹³  →  pipeline Python reproductible et 3× plus rapide que R.",
            size=12, color=TXT)

add_footer(s, 10)

# ═══════════════════════ SAVE ═══════════════════════
prs.save(OUTPATH)
print(f"✓ Wrote {OUTPATH}")
print(f"  Slides: {len(prs.slides)}")
print(f"  Size:   {OUTPATH.stat().st_size / 1024:.1f} KB")
